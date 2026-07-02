from __future__ import annotations

import zipfile
from pathlib import Path

from driftmd.records import StepRecord, append_record, phase_names, read_json, utc_now


def _workflow_summary(root: Path) -> str:
    phases = phase_names(root)
    if {"prepare", "simulate", "analyze"}.issubset(phases):
        return "This report summarizes a complete prepare, simulate, and analyze workflow."
    if "analyze" in phases and "simulate" not in phases:
        return (
            "This report was generated from an existing trajectory. "
            "Preparation and simulation were not run in this workflow."
        )
    return "This report summarizes the workflow phases recorded in the run manifest."


def _write_slides_outline(root: Path, out: Path, title: str) -> None:
    phases = phase_names(root)
    lines = [
        f"# {title}",
        "",
        "## Overview",
        f"- {_workflow_summary(root)}",
    ]
    if "prepare" in phases:
        lines.extend(["", "## Preparation", "- See `prepare/prepare_manifest.json`."])
    if "simulate" in phases:
        lines.extend(["", "## Simulation", "- See `simulate/simulation_manifest.json`."])
    if "analyze" in phases:
        lines.extend(["", "## Analysis", "- See `analysis/analysis_manifest.json`."])
    (out / "slides_outline.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


def _write_pptx_if_available(root: Path, out: Path, title: str) -> list[str]:
    try:
        from pptx import Presentation
    except ImportError:
        return []
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = "DriftMD Workbench"
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Workflow"
    slide.placeholders[1].text = _workflow_summary(root)
    pptx = out / "slides.pptx"
    prs.save(str(pptx))
    return ["slides.pptx"]


def _write_bundle(root: Path, out: Path) -> None:
    bundle = out / "workflow_bundle.zip"
    with zipfile.ZipFile(bundle, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for path in sorted(root.rglob("*")):
            if not path.is_file() or path == bundle:
                continue
            if "__pycache__" in path.parts or path.suffix in {".pyc", ".tmp"}:
                continue
            zf.write(path, path.relative_to(root).as_posix())


def build_report(root: Path, title: str = "DriftMD Workflow Report") -> StepRecord:
    started = utc_now()
    out = root / "report"
    out.mkdir(parents=True, exist_ok=True)
    analysis = read_json(root / "analysis" / "analysis_manifest.json", {})
    lines = [
        f"# {title}",
        "",
        "## Summary",
        "",
        _workflow_summary(root),
        "",
        "## Analysis outputs",
        "",
    ]
    if analysis:
        for artifact in analysis.get("artifacts", []):
            lines.append(f"- `analysis/{artifact}`")
    else:
        lines.append("No analysis manifest was found.")
    (out / "report.md").write_text("\n".join(lines) + "\n", encoding="utf-8")
    _write_slides_outline(root, out, title)
    pptx_artifacts = _write_pptx_if_available(root, out, title)
    _write_bundle(root, out)
    artifacts = ["report.md", "slides_outline.md", *pptx_artifacts, "workflow_bundle.zip"]
    record = StepRecord(
        name="report",
        status="ok",
        started_at=started,
        finished_at=utc_now(),
        output_dir=str(out),
        artifacts=artifacts,
        message="Report artifacts generated.",
    )
    append_record(root, record)
    return record
