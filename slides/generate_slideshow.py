from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DEMOS_DIR = ROOT / "demos"
OUTPUT_PPTX = ROOT / "slides" / "FastMDAnalysis_demo_slides.pptx"

COMMANDS = {
    "rmsd": (
    'fastmda rmsd -traj data/trp_cage.dcd -top data/trp_cage.pdb '
    '--reference-frame 0 --atoms "protein and name CA" --frames 0,-1,10 -o demos/rmsd'
    ),
    "rmsf": (
    'fastmda rmsf -traj data/trp_cage.dcd -top data/trp_cage.pdb '
    '--atoms "protein and name CA" --frames 0,-1,10 -o demos/rmsf'
    ),
    "rg": (
    'fastmda rg -traj data/trp_cage.dcd -top data/trp_cage.pdb '
    '--atoms "protein" --frames 0,-1,10 -o demos/rg'
    ),
    "hbonds": (
    'fastmda hbonds -traj data/trp_cage.dcd -top data/trp_cage.pdb '
    '--atoms "protein" --frames 0,-1,10 -o demos/hbonds'
    ),
    "sasa": (
    'fastmda sasa -traj data/trp_cage.dcd -top data/trp_cage.pdb '
    '--atoms "protein" --probe_radius 0.14 --frames 0,-1,10 -o demos/sasa'
    ),
    "ss": (
    'fastmda ss -traj data/trp_cage.dcd -top data/trp_cage.pdb '
    '--atoms "protein" --frames 0,-1,10 -o demos/ss'
    ),
    "cluster": (
    'fastmda cluster -traj data/trp_cage.dcd -top data/trp_cage.pdb '
    '--atoms "protein and name CA" --methods dbscan kmeans hierarchical --eps 0.40 '
        '--min_samples 8 --n_clusters 4 --frames 0,-1,10 -o demos/cluster'
    ),
    "dimred": (
    'fastmda dimred -traj data/trp_cage.dcd -top data/trp_cage.pdb '
    '--methods pca mds tsne --atoms "protein and name CA" '
        '--frames 0,-1,10 -o demos/dimred'
    ),
}


def add_slide(prs: Presentation, image_path: Path, command: str) -> None:
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    slide.shapes.add_picture(
        str(image_path),
        left=Inches(0.5),
        top=Inches(0.5),
        width=Inches(8.5),
    )

    textbox = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.4), Inches(8.5), Inches(1.5)
    )
    frame = textbox.text_frame
    frame.word_wrap = True

    title_para = frame.paragraphs[0]
    title_run = title_para.add_run()
    title_run.text = f"{image_path.parent.name.upper()} – {image_path.name}"
    title_run.font.size = Pt(20)
    title_run.font.bold = True

    code_para = frame.add_paragraph()
    code_run = code_para.add_run()
    code_run.text = command
    code_run.font.name = "Consolas"
    code_run.font.size = Pt(14)


def main() -> None:
    prs = Presentation()

    for analysis, command in COMMANDS.items():
        image_folder = DEMOS_DIR / analysis
        if not image_folder.exists():
            continue
        for image_path in sorted(image_folder.glob("*.png")):
            add_slide(prs, image_path, command)

    OUTPUT_PPTX.parent.mkdir(exist_ok=True)
    prs.save(OUTPUT_PPTX)
    print(f"Saved slideshow to {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
