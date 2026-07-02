"""Tests for the report-phase consumption of analysis-phase outputs.

These tests verify that after a real analysis run, the report phase:
  1. Embeds analysis figures into report.md with correct relative paths
  2. Shows per-analysis options (from the options.json files) in the document
  3. Generates image slides in the .pptx for each analysis figure
  4. Fans out multi-method analyses (cluster, dimred) into separate slides
"""

from __future__ import annotations

import json
import zipfile
from pathlib import Path

import mdtraj as md
import numpy as np
import pytest

from fastmdxplora import FastMDXplora
from fastmdxplora.cli.main import main as cli_main


# ---------------------------------------------------------------------------
# Helper: build a real trajectory + run a real analysis phase
# ---------------------------------------------------------------------------
def _make_traj_files(tmp_path: Path, *, n_residues: int = 5, n_frames: int = 25) -> Path:
    """Save a compact-globule trajectory under simulation/ so analysis can find it.

    Returns the project root (the directory containing simulation/).
    """
    rng = np.random.RandomState(0)
    top = md.Topology()
    chain = top.add_chain()
    for i in range(n_residues):
        res = top.add_residue("ALA", chain, resSeq=i + 1)
        for nm in ("N", "CA", "C", "O", "CB"):
            el = {
                "N": md.element.nitrogen,
                "CA": md.element.carbon,
                "C": md.element.carbon,
                "O": md.element.oxygen,
                "CB": md.element.carbon,
            }[nm]
            top.add_atom(nm, el, res)
    top.create_standard_bonds()

    base = rng.uniform(-0.2, 0.2, size=(top.n_atoms, 3))
    xyz = np.tile(base[None, :, :], (n_frames, 1, 1)) + rng.normal(
        scale=0.02, size=(n_frames, top.n_atoms, 3)
    )
    traj = md.Trajectory(
        xyz=xyz.astype(np.float32), topology=top, time=np.arange(n_frames) * 20.0
    )

    sim = tmp_path / "simulation"
    sim.mkdir(parents=True)
    traj[0].save_pdb(str(sim / "topology.pdb"))
    traj.save_dcd(str(sim / "production.dcd"))
    return tmp_path


@pytest.fixture
def project_with_analysis(tmp_path: Path) -> Path:
    """Project directory after a real analysis phase has run on RMSD + Rg."""
    root = _make_traj_files(tmp_path / "proj")
    fmdx = FastMDXplora(
        system=str(root / "simulation" / "topology.pdb"),
        output_dir=root,
    )
    fmdx.explore(
        include=["analysis"],
        options={"analysis": {"include": ["rmsd", "rg"]}},
    )
    return root


@pytest.fixture
def project_with_multi_method(tmp_path: Path) -> Path:
    """Project after analysis with multi-method cluster + dimred."""
    root = _make_traj_files(tmp_path / "proj", n_residues=6, n_frames=30)
    fmdx = FastMDXplora(
        system=str(root / "simulation" / "topology.pdb"),
        output_dir=root,
    )
    fmdx.explore(
        include=["analysis"],
        options={
            "analysis": {
                "include": ["cluster", "dimred"],
                "options": {
                    "cluster": {
                        "methods": ["kmeans", "hierarchical", "dbscan"],
                        "n_clusters": 3,
                    },
                    "dimred": {"methods": ["pca", "tsne"]},
                },
            }
        },
    )
    return root


# ===========================================================================
# Markdown report
# ===========================================================================
class TestReportDocument:
    def test_report_embeds_analysis_figures(self, project_with_analysis: Path):
        """The Markdown report should reference each analysis's PNG file."""
        fmdx = FastMDXplora(
            system=str(project_with_analysis / "simulation" / "topology.pdb"),
            output_dir=project_with_analysis,
        )
        fmdx.report()
        text = (project_with_analysis / "report" / "report.md").read_text(encoding="utf-8")
        # Each analysis's figure should be referenced
        assert "analysis_summary.png" in text
        assert "analysis/rmsd/rmsd.png" in text
        assert "analysis/rg/rg.png" in text
        # No "deferred" / placeholder language
        assert "deferred" not in text.lower()
        assert "wired in" not in text.lower()

    def test_report_shows_actual_options(self, project_with_analysis: Path):
        """Options from options.json should appear in the report parameters block."""
        fmdx = FastMDXplora(
            system=str(project_with_analysis / "simulation" / "topology.pdb"),
            output_dir=project_with_analysis,
        )
        fmdx.report()
        text = (project_with_analysis / "report" / "report.md").read_text(encoding="utf-8")
        # RMSD options should appear
        assert "`align`" in text
        assert "`ref`" in text
        # And the actual values, not the placeholder
        assert "Run with default options" not in text

    def test_report_includes_trajectory_metadata(self, project_with_analysis: Path):
        """The report should mention how many frames/residues were analyzed."""
        fmdx = FastMDXplora(
            system=str(project_with_analysis / "simulation" / "topology.pdb"),
            output_dir=project_with_analysis,
        )
        fmdx.report()
        text = (project_with_analysis / "report" / "report.md").read_text(encoding="utf-8")
        # We built a 25-frame, 5-residue trajectory
        assert "25 frames" in text
        assert "5 residues" in text

    def test_report_includes_all_figures_for_multi_method(
        self, project_with_multi_method: Path
    ):
        """For cluster/dimred, every method's PNG should be referenced."""
        fmdx = FastMDXplora(
            system=str(project_with_multi_method / "simulation" / "topology.pdb"),
            output_dir=project_with_multi_method,
        )
        fmdx.report()
        text = (project_with_multi_method / "report" / "report.md").read_text(encoding="utf-8")
        # cluster has two methods
        assert "cluster/cluster_kmeans.png" in text
        assert "cluster/cluster_kmeans_counts.png" in text
        assert "cluster/cluster_hierarchical.png" in text
        assert "cluster/cluster_hierarchical_counts.png" in text
        assert "cluster/cluster_hierarchical_dendrogram.png" in text
        # dimred has two methods
        assert "dimred/dimred_pca.png" in text
        assert "dimred/dimred_tsne.png" in text

    def test_report_writes_combined_analysis_summary(
        self, project_with_multi_method: Path
    ):
        fmdx = FastMDXplora(
            system=str(project_with_multi_method / "simulation" / "topology.pdb"),
            output_dir=project_with_multi_method,
        )
        fmdx.report()

        summary = project_with_multi_method / "report" / "analysis_summary.png"
        summary_manifest = (
            project_with_multi_method / "report" / "analysis_summary_manifest.json"
        )
        assert summary.is_file()
        assert summary_manifest.is_file()
        manifest = json.loads(summary_manifest.read_text(encoding="utf-8"))
        included_sources = {item["source"] for item in manifest["included"]}
        assert "analysis/cluster/cluster_kmeans_counts.png" in included_sources
        assert "analysis/cluster/cluster_hierarchical_dendrogram.png" in included_sources
        assert "analysis/dimred/dimred_pca.png" in included_sources


# ===========================================================================
# PPTX slide deck
# ===========================================================================
class TestSlideDeck:
    def test_pptx_embeds_analysis_figures(self, project_with_analysis: Path):
        """Each analysis figure should appear as an image on its own slide."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        fmdx = FastMDXplora(
            system=str(project_with_analysis / "simulation" / "topology.pdb"),
            output_dir=project_with_analysis,
        )
        fmdx.report()
        prs = Presentation(str(project_with_analysis / "report" / "slides.pptx"))

        # Count image slides
        image_slides = [
            s for s in prs.slides
            if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in s.shapes)
        ]
        # We ran RMSD + Rg, plus the combined summary slide.
        assert len(image_slides) == 3

    def test_pptx_image_slides_titled_by_analysis(self, project_with_analysis: Path):
        """Image-slide titles should be readable, not giant all-caps IDs."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        fmdx = FastMDXplora(
            system=str(project_with_analysis / "simulation" / "topology.pdb"),
            output_dir=project_with_analysis,
        )
        fmdx.report()
        prs = Presentation(str(project_with_analysis / "report" / "slides.pptx"))

        titles = {
            slide.shapes.title.text
            for slide in prs.slides
            if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)
        }
        assert "Analysis summary" in titles
        assert "RMSD over frames" in titles
        assert "Radius of gyration" in titles
        assert "RG" not in titles

    def test_pptx_fans_out_multi_method_analyses(
        self, project_with_multi_method: Path
    ):
        """Multi-method analyses should produce individual restored plot slides."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        fmdx = FastMDXplora(
            system=str(project_with_multi_method / "simulation" / "topology.pdb"),
            output_dir=project_with_multi_method,
        )
        fmdx.report()
        prs = Presentation(str(project_with_multi_method / "report" / "slides.pptx"))

        image_slides = [
            s for s in prs.slides
            if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in s.shapes)
        ]
        # cluster: timelines + count plots + dendrogram, dimred: 2 methods,
        # plus the combined summary slide.
        assert len(image_slides) >= 8

    def test_pptx_multi_method_subtitle_records_method(
        self, project_with_multi_method: Path
    ):
        """Multi-method image slides should have a subtitle indicating the method."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        fmdx = FastMDXplora(
            system=str(project_with_multi_method / "simulation" / "topology.pdb"),
            output_dir=project_with_multi_method,
        )
        fmdx.report()
        prs = Presentation(str(project_with_multi_method / "report" / "slides.pptx"))

        # Find all text on image-bearing slides
        all_text = []
        for slide in prs.slides:
            if not any(
                shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes
            ):
                continue
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text.append(shape.text_frame.text)
        joined = "\n".join(all_text).lower()
        # The subtitle text should include method names somewhere
        for method in ("kmeans", "hierarchical", "pca", "tsne"):
            assert method in joined, f"method {method} missing from slide text"
        assert "analysis summary" in joined


# ===========================================================================
# Static HTML dashboard
# ===========================================================================
class TestDashboard:
    def test_dashboard_is_generated_and_links_outputs(self, project_with_analysis: Path):
        rmsd_png = project_with_analysis / "analysis" / "rmsd" / "rmsd.png"
        rg_png = project_with_analysis / "analysis" / "rg" / "rg.png"
        original_rmsd = rmsd_png.read_bytes()
        original_rg = rg_png.read_bytes()
        fmdx = FastMDXplora(
            system=str(project_with_analysis / "simulation" / "topology.pdb"),
            output_dir=project_with_analysis,
        )
        result = fmdx.report()

        assert result.status == "ok"
        dashboard = project_with_analysis / "report" / "dashboard.html"
        assert dashboard.is_file()
        text = dashboard.read_text(encoding="utf-8")

        assets = project_with_analysis / "report" / "dashboard_assets"
        assert (assets / "rmsd_dashboard.png").is_file()
        assert (assets / "rg_dashboard.png").is_file()
        assert rmsd_png.read_bytes() == original_rmsd
        assert rg_png.read_bytes() == original_rg

        assert '<aside class="sidebar">' in text
        assert "Run Progress" in text
        assert "Live Simulation" in text
        assert "Live simulation telemetry was not recorded for this run" in text
        assert "Top Metrics" in text
        assert "Recent Outputs" in text
        assert "Quick Actions" in text
        assert 'class="plot-frame"' in text
        assert 'class="plot-grid"' in text
        assert "plot-card card-md" in text
        assert "card-sm" in text
        assert "card-lg" in text
        assert "card-wide" in text
        assert "resize-handle" in text
        assert 'title="Drag to resize"' in text
        assert 'data-card-size="sm"' in text
        assert 'data-card-size="md"' in text
        assert 'data-card-size="lg"' in text
        assert 'data-card-size="wide"' in text
        assert "Reset layout" in text
        assert "data-reset-layout" in text
        assert "grid-auto-flow: dense" in text
        assert "grid-auto-rows: 8px" in text
        assert "grid-column: span var(--col-span, 1)" in text
        assert "grid-row: span var(--row-span, 20)" in text
        assert "pointerdown" in text
        assert "pointermove" in text
        assert "setPointerCapture" in text
        assert "localStorage" in text
        assert "ResizeObserver" not in text
        assert "resize: both" not in text
        assert ".panels" not in text
        assert "repeat(auto-fill, minmax(280px, 1fr))" in text
        assert "min-height: 180px" in text
        assert "overflow: hidden" in text
        assert "max-height: 100%" in text
        assert "object-fit: contain" in text
        assert "plot-card large" not in text
        assert ".plot-card.card-lg { --col-span: 2; --row-span: 38; }" in text
        assert ".plot-card.card-wide { --col-span: 2; --row-span: 24; }" in text
        assert "dashboard view" in text
        assert "dashboard_assets/rmsd_dashboard.png" in text
        assert "dashboard_assets/rg_dashboard.png" in text
        assert "Open Markdown Report" in text
        assert "Open Analysis Manifest" in text
        assert "output-list" in text
        assert "outputs-extra" in text
        assert "repeat(auto-fit, minmax(220px, 1fr))" in text
        assert "action-title" in text
        assert "action-subtitle" in text
        assert "artifact-path" in text
        assert "Std. Dev." in text
        assert "original: <a href=\"../analysis/rmsd/rmsd.png\"" in text
        assert "original: <a href=\"../analysis/rg/rg.png\"" in text
        assert "../analysis/analysis_manifest.json" in text
        assert "report.md" in text
        assert "slides.pptx" in text
        assert "project_bundle.zip" in text
        assert "dashboard.html" in text
        assert "TODO" not in text
        assert "fake" not in text.lower()
        assert "dummy" not in text.lower()
        with zipfile.ZipFile(project_with_analysis / "report" / "project_bundle.zip") as zf:
            names = set(zf.namelist())
        assert "report/dashboard_assets/rmsd_dashboard.png" in names
        assert "report/dashboard_assets/rg_dashboard.png" in names

    def test_dashboard_analysis_only_phase_aware_text(self, tmp_path: Path):
        root = tmp_path / "analysis_only_dashboard"
        analysis = root / "analysis"
        rmsd_dir = analysis / "rmsd"
        rg_dir = analysis / "rg"
        rmsd_dir.mkdir(parents=True)
        rg_dir.mkdir()
        (rmsd_dir / "rmsd.png").write_bytes(b"png")
        (rg_dir / "rg.png").write_bytes(b"png")
        (root / "manifest.json").write_text(
            json.dumps(
                {
                    "system": "1L2Y",
                    "phases": [
                        {"name": "analysis", "status": "ok"},
                        {"name": "report", "status": "ok"},
                    ],
                }
            ),
            encoding="utf-8",
        )
        (analysis / "analysis_manifest.json").write_text(
            json.dumps(
                {
                    "plan": ["rmsd", "rg"],
                    "n_frames": 10,
                    "n_atoms": 100,
                    "results": {"rmsd": {"status": "ok"}, "rg": {"status": "ok"}},
                }
            ),
            encoding="utf-8",
        )

        fmdx = FastMDXplora(system="1L2Y", output_dir=root)
        result = fmdx.report(title="Analysis-only dashboard", slides=False, bundle=False)

        assert result.status == "ok"
        text = (root / "report" / "dashboard.html").read_text(encoding="utf-8")
        assert (
            "Analysis/report workflow from existing trajectory. Setup and "
            "simulation were not run in this workflow."
        ) in text
        assert "Existing trajectory analysis" in text
        assert "<span>Setup</span><span class=\"phase-detail\">Not run</span>" in text
        assert "<span>Simulation</span><span class=\"phase-detail\">Not run</span>" in text
        assert "<span>Analysis</span><span class=\"phase-detail\">Completed</span>" in text
        assert "<span>Report</span><span class=\"phase-detail\">Completed</span>" in text
        assert 'class="plot-frame"' in text
        assert "artifact fallback" in text
        assert "dashboard_assets" not in text
        assert "Quick Actions" in text
        assert "Recent Outputs" in text
        assert "Simulation time" not in text
        assert "Temperature" not in text
        assert "Production MD completed" not in text
        assert "../analysis/rmsd/rmsd.png" in text
        assert "../analysis/rg/rg.png" in text
        assert "project_bundle.zip" not in text

    def test_dashboard_single_artifact_sections_use_normal_plot_grid(
        self, tmp_path: Path
    ):
        root = tmp_path / "single_artifact_sections"
        artifacts = [
            "analysis/sasa/sasa.png",
            "analysis/ss/ss.png",
            "analysis/dimred/dimred_pca.png",
        ]
        png = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
            b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
            b"\x00\x00\x00\x0cIDATx\x9cc```\x00\x00\x00\x04\x00\x01"
            b"\xf6\x178U\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        for rel in artifacts:
            path = root / rel
            path.parent.mkdir(parents=True, exist_ok=True)
            path.write_bytes(png)
        (root / "manifest.json").write_text(
            json.dumps(
                {
                    "system": "1L2Y",
                    "phases": [
                        {"name": "analysis", "status": "ok"},
                        {"name": "report", "status": "ok"},
                    ],
                }
            ),
            encoding="utf-8",
        )
        (root / "analysis" / "analysis_manifest.json").write_text(
            json.dumps(
                {
                    "plan": ["sasa", "ss", "dimred"],
                    "results": {
                        "sasa": {"status": "ok"},
                        "ss": {"status": "ok"},
                        "dimred": {"status": "ok"},
                    },
                }
            ),
            encoding="utf-8",
        )

        fmdx = FastMDXplora(system="1L2Y", output_dir=root)
        result = fmdx.report(document=False, slides=False, bundle=False)

        assert result.status == "ok"
        text = (root / "report" / "dashboard.html").read_text(encoding="utf-8")
        assert text.count('class="plot-grid"') == 1
        assert "Additional Analysis" in text
        assert 'id="additional-analysis"' in text
        assert 'id="sasa-section"' not in text
        assert 'id="secondary-structure-section"' not in text
        assert 'id="dimensionality-reduction"' not in text
        assert "Total SASA" in text
        assert "Secondary structure" in text
        assert "PCA" in text
        assert text.count('class="plot-card fallback card-md"') == len(artifacts)
        assert text.count('class="resize-handle"') == len(artifacts)
        assert text.count('class="plot-frame"') == len(artifacts)
        assert "plot-card large" not in text
        assert ".plot-card.card-lg { --col-span: 2; --row-span: 38; }" in text
        assert ".plot-card.card-wide { --col-span: 2; --row-span: 24; }" in text
        assert "section-secondary-structure-section { --plot-min" not in text
        assert "section-sasa-section { --plot-min" not in text
        assert "artifact fallback" in text
        assert "data-card-key=" in text

    def test_dashboard_multi_method_dark_assets(self, project_with_multi_method: Path):
        original_pngs = {
            path: path.read_bytes()
            for path in (
                project_with_multi_method / "analysis" / "cluster" / "cluster_kmeans.png",
                project_with_multi_method
                / "analysis"
                / "cluster"
                / "cluster_hierarchical.png",
                project_with_multi_method / "analysis" / "dimred" / "dimred_pca.png",
            )
        }
        fmdx = FastMDXplora(
            system=str(project_with_multi_method / "simulation" / "topology.pdb"),
            output_dir=project_with_multi_method,
        )
        result = fmdx.report()

        assert result.status == "ok"
        text = (
            project_with_multi_method / "report" / "dashboard.html"
        ).read_text(encoding="utf-8")
        assets = project_with_multi_method / "report" / "dashboard_assets"
        expected_assets = [
            "pca_dashboard.png",
            "tsne_dashboard.png",
            "kmeans_trajectory_dashboard.png",
            "kmeans_population_dashboard.png",
            "dbscan_trajectory_dashboard.png",
            "dbscan_population_dashboard.png",
            "hierarchical_trajectory_dashboard.png",
            "hierarchical_population_dashboard.png",
            "hierarchical_dendrogram_dashboard.png",
        ]
        for name in expected_assets:
            assert (assets / name).is_file()
            assert f"dashboard_assets/{name}" in text
        assert "Hierarchical dendrogram" in text
        assert "dashboard_assets/hierarchical_dendrogram_dashboard.png" in text
        assert "original: <a href=\"../analysis/cluster/cluster_hierarchical_dendrogram.png\"" in text
        assert "dashboard_assets/dbscan_trajectory_dashboard.png" in text
        assert "dashboard_assets/dbscan_population_dashboard.png" in text
        assert "dashboard view" in text
        with zipfile.ZipFile(
            project_with_multi_method / "report" / "project_bundle.zip"
        ) as zf:
            names = set(zf.namelist())
        for name in expected_assets:
            assert f"report/dashboard_assets/{name}" in names
        for path, data in original_pngs.items():
            assert path.read_bytes() == data

    def test_dashboard_discovers_all_analysis_image_artifacts(self, tmp_path: Path):
        root = tmp_path / "artifact_complete"
        artifact_names = [
            "analysis/sasa/total_sasa.png",
            "analysis/sasa/residue_sasa.png",
            "analysis/sasa/average_residue_sasa.png",
            "analysis/dimred/dimred_pca.png",
            "analysis/dimred/dimred_mds.png",
            "analysis/dimred/dimred_tsne.png",
            "analysis/cluster/dbscan_pop.png",
            "analysis/cluster/dbscan_traj_hist.png",
            "analysis/cluster/dbscan_traj_scatter.png",
            "analysis/cluster/dbscan_distance_matrix.png",
            "analysis/cluster/kmeans_pop.png",
            "analysis/cluster/kmeans_traj_hist.png",
            "analysis/cluster/kmeans_traj_scatter.png",
            "analysis/cluster/hierarchical_pop.png",
            "analysis/cluster/hierarchical_traj_hist.png",
            "analysis/cluster/hierarchical_traj_scatter.png",
            "analysis/cluster/hierarchical_dendrogram.png",
        ]
        png = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
            b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
            b"\x00\x00\x00\x0cIDATx\x9cc```\x00\x00\x00\x04\x00\x01"
            b"\xf6\x178U\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        for rel in artifact_names:
            path = root / rel
            path.parent.mkdir(parents=True, exist_ok=True)
            path.write_bytes(png)
        (root / "manifest.json").write_text(
            json.dumps(
                {
                    "system": "1L2Y",
                    "phases": [
                        {"name": "analysis", "status": "ok"},
                        {"name": "report", "status": "ok"},
                    ],
                }
            ),
            encoding="utf-8",
        )
        (root / "analysis" / "analysis_manifest.json").write_text(
            json.dumps(
                {
                    "plan": ["sasa", "dimred", "cluster"],
                    "results": {
                        "sasa": {"status": "ok"},
                        "dimred": {"status": "ok"},
                        "cluster": {"status": "ok"},
                    },
                }
            ),
            encoding="utf-8",
        )

        fmdx = FastMDXplora(system="1L2Y", output_dir=root)
        result = fmdx.report(document=False, slides=False, bundle=False)

        assert result.status == "ok"
        text = (root / "report" / "dashboard.html").read_text(encoding="utf-8")
        assert text.count('<div class="plot-frame">') == len(artifact_names)
        assert text.count('class="plot-card fallback card-md"') == len(artifact_names)
        assert text.count('class="resize-handle"') == len(artifact_names)
        assert text.count('<span class="tag">artifact fallback</span>') == len(artifact_names)
        assert "artifact fallback" in text
        assert "dashboard view" not in text
        assert "output-list" in text
        assert "outputs-extra" in text
        assert "Show all outputs" in text
        assert "plot-card large" not in text
        assert ".plot-card.card-lg { --col-span: 2; --row-span: 38; }" in text
        assert ".plot-card.card-wide { --col-span: 2; --row-span: 24; }" in text
        for rel in artifact_names:
            assert f"../{rel}" in text
        for label in (
            "Total SASA",
            "Per-residue SASA heatmap",
            "Average per-residue SASA",
            "PCA",
            "MDS",
            "t-SNE",
            "DBSCAN population plot",
            "DBSCAN trajectory histogram",
            "DBSCAN trajectory scatter",
            "DBSCAN distance matrix",
            "KMeans population plot",
            "KMeans trajectory histogram",
            "KMeans trajectory scatter",
            "Hierarchical population plot",
            "Hierarchical trajectory histogram",
            "Hierarchical trajectory scatter",
            "Hierarchical dendrogram",
        ):
            assert label in text
        assert "Core Metrics" not in text
        assert "SASA" in text
        assert "Dimensionality Reduction" in text
        assert "Clustering" in text


# ===========================================================================
# End-to-end with no analyses (deferred case)
# ===========================================================================
def test_report_handles_missing_analyses_gracefully(tmp_path: Path):
    """When no analysis has run, the report should still build a coherent deck."""
    fmdx = FastMDXplora(
        system="1L2Y",  # no trajectory at all
        output_dir=tmp_path,
    )
    result = fmdx.report()
    assert result.status == "ok"
    # The deck should exist; image-slide count is zero because no analyses ran
    from pptx import Presentation

    prs = Presentation(str(tmp_path / "report" / "slides.pptx"))
    assert len(prs.slides) >= 1


def test_deferred_analysis_message_is_current_and_actionable(tmp_path: Path):
    fmdx = FastMDXplora(system="1L2Y", output_dir=tmp_path)
    result = fmdx.analyze()

    assert result.status == "ok"
    manifest = json.loads(
        (tmp_path / "analysis" / "analysis_manifest.json").read_text(encoding="utf-8")
    )
    assert manifest["status"] == "deferred"
    assert "v0.2" not in manifest["note"]
    assert "Run the simulation phase first" in manifest["note"]


# ===========================================================================
# Robustness for report-only mode and generated artifacts
# ===========================================================================
def test_report_only_cli_escapes_weird_markdown_text(tmp_path: Path):
    """Report-only mode should not let titles/messages reshape Markdown."""
    root = tmp_path / "existing run with spaces"
    analysis = root / "analysis"
    analysis.mkdir(parents=True)
    weird_system = "sys|`tick` [brackets]\n# injected\tpath"
    weird_title = "Study | `ticks` [link](x)\n# injected"
    weird_author = "Author | # name\nnext"
    weird_message = "bad | `tick`\n# injected failure"

    (root / "manifest.json").write_text(
        json.dumps({"system": weird_system}),
        encoding="utf-8",
    )
    (analysis / "analysis_manifest.json").write_text(
        json.dumps(
            {
                "plan": ["rmsd"],
                "results": {"rmsd": {"status": "error", "message": weird_message}},
            }
        ),
        encoding="utf-8",
    )

    rc = cli_main(
        [
            "report",
            "--output",
            str(root),
            "--title",
            weird_title,
            "--author",
            weird_author,
            "--no-slides",
            "--no-bundle",
        ]
    )

    assert rc == 0
    text = (root / "report" / "report.md").read_text(encoding="utf-8")
    first_line = text.splitlines()[0]
    assert first_line.startswith("# Study")
    assert "\\|" in first_line
    assert "\\`ticks\\`" in first_line
    assert "\n# injected" not in text
    assert "bad \\| \\`tick\\` \\# injected failure" in text


def test_report_phase_manifest_artifacts_exist(tmp_path: Path):
    fmdx = FastMDXplora(system="sys|odd", output_dir=tmp_path / "run")
    results = fmdx.explore(
        include=["report"],
        options={"report": {"title": "Odd | title", "slides": False, "bundle": False}},
    )

    report_phase = results[0].phase("report")
    assert report_phase is not None
    manifest = json.loads((fmdx.output_dir / "manifest.json").read_text(encoding="utf-8"))
    report_record = next(p for p in manifest["phases"] if p["name"] == "report")
    assert report_record["artifacts"] == ["report.md", "dashboard.html"]
    for artifact in report_record["artifacts"]:
        assert (Path(report_record["output_dir"]) / artifact).is_file()


def test_report_bundle_excludes_cache_and_temp_files(tmp_path: Path):
    root = tmp_path / "run"
    (root / "analysis").mkdir(parents=True)
    (root / "__pycache__").mkdir()
    (root / ".cache").mkdir()
    (root / "__pycache__" / "module.cpython-310.pyc").write_bytes(b"cached")
    (root / ".cache" / "plot.tmp").write_text("cache", encoding="utf-8")
    (root / "scratch.tmp").write_text("temp", encoding="utf-8")
    (root / "keep.dat").write_text("artifact", encoding="utf-8")

    fmdx = FastMDXplora(system="1L2Y", output_dir=root)
    result = fmdx.report(document=False, slides=False, bundle=True)

    assert result.status == "ok"
    bundle = root / "report" / "project_bundle.zip"
    with zipfile.ZipFile(bundle) as zf:
        names = set(zf.namelist())
    assert "keep.dat" in names
    assert "report/dashboard.html" in names
    assert "scratch.tmp" not in names
    assert "__pycache__/module.cpython-310.pyc" not in names
    assert ".cache/plot.tmp" not in names
    assert "report/project_bundle.zip" not in names


def test_powerpoint_handles_weird_title_and_system(tmp_path: Path):
    from pptx import Presentation

    title = "T" * 800 + " | `tick`\n# heading"
    system = "system with spaces | [id]\nnext"
    fmdx = FastMDXplora(system=system, output_dir=tmp_path / "run")
    result = fmdx.report(title=title, bundle=False)

    assert result.status == "ok"
    pptx = fmdx.output_dir / "report" / "slides.pptx"
    assert pptx.is_file()
    prs = Presentation(str(pptx))
    assert len(prs.slides) >= 1
    assert "\n# heading" not in prs.slides[0].shapes.title.text


def test_analysis_report_only_wording_uses_existing_trajectory(tmp_path: Path):
    from pptx import Presentation

    root = tmp_path / "analysis_only"
    analysis = root / "analysis"
    analysis.mkdir(parents=True)
    (root / "manifest.json").write_text(
        json.dumps(
            {
                "system": "1L2Y",
                "phases": [
                    {"name": "analysis", "status": "ok"},
                    {"name": "report", "status": "ok"},
                ],
            }
        ),
        encoding="utf-8",
    )
    (analysis / "analysis_manifest.json").write_text(
        json.dumps(
            {
                "plan": ["rmsd"],
                "n_frames": 10,
                "n_residues": 20,
                "results": {"rmsd": {"status": "ok"}},
            }
        ),
        encoding="utf-8",
    )
    rmsd_dir = analysis / "rmsd"
    rmsd_dir.mkdir()
    (rmsd_dir / "options.json").write_text(
        json.dumps({"selection": "name CA", "options": {"align": True}}),
        encoding="utf-8",
    )

    fmdx = FastMDXplora(system="1L2Y", output_dir=root)
    result = fmdx.report(title="Analysis-only Trp-cage report", bundle=False)

    assert result.status == "ok"
    report = (root / "report" / "report.md").read_text(encoding="utf-8")
    assert "end-to-end molecular dynamics study" not in report
    assert "Simulation parameters were not recorded for this run" not in report
    assert "This report was generated from an existing trajectory" in report
    assert "Setup and simulation were not run in this workflow" in report
    assert "Simulation was not run in this workflow" in report

    outline = (root / "report" / "slides_outline.md").read_text(encoding="utf-8")
    assert "Analysis/report workflow from an existing trajectory" in outline
    assert "Setup and simulation were not run in this workflow" in outline
    assert "## 2. Setup" not in outline
    assert "## 3. Simulation" not in outline

    prs = Presentation(str(root / "report" / "slides.pptx"))
    all_text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Analysis/report workflow from an existing trajectory" in all_text
    assert "Setup and simulation were not run in this workflow" in all_text


def test_full_pipeline_report_keeps_end_to_end_wording(tmp_path: Path):
    root = tmp_path / "full_pipeline"
    (root / "setup").mkdir(parents=True)
    (root / "simulation").mkdir()
    (root / "analysis").mkdir()
    (root / "manifest.json").write_text(
        json.dumps(
            {
                "system": "1L2Y",
                "phases": [
                    {"name": "setup", "status": "ok"},
                    {"name": "simulation", "status": "ok"},
                    {"name": "analysis", "status": "ok"},
                    {"name": "report", "status": "ok"},
                ],
            }
        ),
        encoding="utf-8",
    )
    (root / "setup" / "setup_parameters.json").write_text(
        json.dumps({"parameters": {"ph": 7.0}}),
        encoding="utf-8",
    )
    (root / "simulation" / "simulation_parameters.json").write_text(
        json.dumps({"parameters": {"duration_ns": 0.001}}),
        encoding="utf-8",
    )
    (root / "analysis" / "analysis_manifest.json").write_text(
        json.dumps({"plan": [], "results": {}}),
        encoding="utf-8",
    )

    fmdx = FastMDXplora(system="1L2Y", output_dir=root)
    result = fmdx.report(title="Full pipeline report", slides=False, bundle=False)

    assert result.status == "ok"
    report = (root / "report" / "report.md").read_text(encoding="utf-8")
    assert "end-to-end molecular dynamics study" in report
    assert "Production MD was performed" in report
    assert "Setup and simulation were not run in this workflow" not in report


def test_region_highlight_report_from_existing_rmsf_outputs(tmp_path: Path):
    from pptx import Presentation

    root = tmp_path / "region_run"
    rmsf_dir = root / "analysis" / "rmsf"
    rmsf_dir.mkdir(parents=True)
    np.savetxt(
        rmsf_dir / "rmsf.dat",
        np.column_stack([np.arange(1, 9), np.linspace(0.02, 0.12, 8)]),
    )
    (root / "manifest.json").write_text(
        json.dumps(
            {
                "system": "example.pdb",
                "phases": [
                    {"name": "analysis", "status": "ok"},
                    {"name": "report", "status": "ok"},
                ],
            }
        ),
        encoding="utf-8",
    )
    (root / "analysis" / "analysis_manifest.json").write_text(
        json.dumps(
            {
                "plan": ["rmsf"],
                "n_frames": 5,
                "n_residues": 8,
                "topology_input": None,
                "results": {"rmsf": {"status": "ok"}},
            }
        ),
        encoding="utf-8",
    )

    fmdx = FastMDXplora(system="example.pdb", output_dir=root)
    result = fmdx.report(
        bundle=False,
        region_highlights=[
            {"label": "example region 1", "start": 2, "end": 4, "color": "#4E79A7"},
            {"label": "example region 2", "start": 6, "end": 7, "color": "#F28E2B"},
        ],
    )

    assert result.status == "ok"
    assert (root / "analysis" / "rmsf" / "rmsf_region_highlights.png").is_file()
    assert (root / "report" / "region_highlight_summary.png").is_file()
    manifest = json.loads(
        (root / "report" / "region_highlight_manifest.json").read_text(encoding="utf-8")
    )
    assert manifest["status"] == "ok"
    assert manifest["skipped"][0]["artifact"] == "structure_region_highlights.png"
    reason = manifest["skipped"][0]["reason"]
    assert "PyMOL" in reason or "topology" in reason

    report = (root / "report" / "report.md").read_text(encoding="utf-8")
    assert "Region Highlight Figure" in report
    assert "region_highlight_summary.png" in report
    outline = (root / "report" / "slides_outline.md").read_text(encoding="utf-8")
    assert "region_highlight_summary.png" in outline
    prs = Presentation(str(root / "report" / "slides.pptx"))
    titles = [s.shapes.title.text for s in prs.slides if s.shapes.title is not None]
    assert "Region highlights" in titles


def test_region_highlight_invalid_range_records_error(tmp_path: Path):
    root = tmp_path / "bad_region"
    rmsf_dir = root / "analysis" / "rmsf"
    rmsf_dir.mkdir(parents=True)
    np.savetxt(rmsf_dir / "rmsf.dat", np.column_stack([np.arange(1, 5), np.ones(4)]))
    (root / "analysis" / "analysis_manifest.json").write_text(
        json.dumps({"plan": ["rmsf"], "results": {"rmsf": {"status": "ok"}}}),
        encoding="utf-8",
    )

    fmdx = FastMDXplora(system="example.pdb", output_dir=root)
    result = fmdx.report(
        document=False,
        slides=False,
        bundle=False,
        region_highlights=[{"label": "bad", "start": 0, "end": 2}],
    )

    assert result.status == "ok"
    manifest = json.loads(
        (root / "report" / "region_highlight_manifest.json").read_text(encoding="utf-8")
    )
    assert manifest["status"] == "error"
    assert "start must be >= 1" in manifest["error"]
    assert not (root / "analysis" / "rmsf" / "rmsf_region_highlights.png").exists()


def test_report_without_region_highlights_has_no_region_artifacts(
    project_with_analysis: Path,
):
    fmdx = FastMDXplora(
        system=str(project_with_analysis / "simulation" / "topology.pdb"),
        output_dir=project_with_analysis,
    )
    fmdx.report(bundle=False)

    assert not (project_with_analysis / "report" / "region_highlight_summary.png").exists()
    text = (project_with_analysis / "report" / "report.md").read_text(encoding="utf-8")
    assert "Region Highlight Figure" not in text


def test_region_highlight_structure_panel_when_pymol_available(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
):
    import subprocess

    from fastmdxplora.report import region_highlights as rh

    def fake_run(cmd, check, text, capture_output, timeout):
        import matplotlib.pyplot as plt

        script_path = Path(cmd[-1])
        script = script_path.read_text(encoding="utf-8")
        out_line = next(line for line in script.splitlines() if line.startswith("png "))
        out_path = Path(out_line.split(" ", 1)[1].split(",", 1)[0])
        plt.imsave(out_path, np.ones((4, 4, 3)))
        return subprocess.CompletedProcess(cmd, 0, "", "")

    monkeypatch.setattr(
        rh,
        "detect_pymol_renderer",
        lambda: rh.PymolRenderer(kind="command", command=("pymol", "-cq")),
    )
    monkeypatch.setattr(rh.subprocess, "run", fake_run)

    root = _make_traj_files(tmp_path / "region_structure", n_residues=8, n_frames=10)
    fmdx = FastMDXplora(
        system=str(root / "simulation" / "topology.pdb"),
        output_dir=root,
    )
    fmdx.explore(
        include=["analysis", "report"],
        options={
            "analysis": {"include": ["rmsf"]},
            "report": {
                "bundle": False,
                "region_highlights": [
                    {"label": "example region", "start": 2, "end": 5}
                ],
            },
        },
    )

    assert (root / "report" / "structure_region_highlights.png").is_file()
    assert (root / "report" / "structure_region_highlights.pml").is_file()
    assert (root / "report" / "region_highlight_summary.png").is_file()
    manifest = json.loads(
        (root / "report" / "region_highlight_manifest.json").read_text(encoding="utf-8")
    )
    assert manifest["renderer"] == "PyMOL"
    assert "report/structure_region_highlights.png" in manifest["artifacts"]
    assert "report/structure_region_highlights.pml" in manifest["artifacts"]
    assert manifest["skipped"] == []


def test_pymol_script_generation_uses_cartoon_and_regions(tmp_path: Path):
    from fastmdxplora.report.region_highlights import (
        RegionHighlight,
        build_pymol_script,
    )

    script = build_pymol_script(
        topology_path=tmp_path / "topology.pdb",
        output_path=tmp_path / "structure.png",
        regions=[
            RegionHighlight("example A", 2, 5, "#4E79A7"),
            RegionHighlight("example B", 7, 9, "orange"),
        ],
    )

    assert "show cartoon, prot" in script
    assert "color gray70, prot" in script
    assert "resi 2-5" in script
    assert "resi 7-9" in script
    assert "ray 1800, 1200" in script
    assert "png " in script
